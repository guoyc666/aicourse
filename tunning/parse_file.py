import time
import os
import fitz
from docx import Document
from pptx import Presentation
import whisper
from .parse_image import understand_image

def extract_text(path: str, file_type: str) -> str:
    """根据文件类型解析文件内容"""
    if file_type == "txt":
        return extract_text_text(path)
    elif file_type == "docx":
        return extract_text_docx(path)
    elif file_type == "pdf":
        return extract_text_pdf(path)
    elif file_type == "pptx":
        return extract_text_pptx(path)
    elif file_type in ["png", "jpg", "jpeg"]:
        return extract_text_image(path)
    elif file_type in ["mp3", "wav", "mp4", "mkv"]:
        return extract_text_audio_vedio(path)
    else:
        raise ValueError(f"不支持的文件类型: {file_type}")


def extract_text_text(path: str) -> str:
    """解析文本文件"""
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def extract_text_docx(path: str) -> str:
    """解析 DOCX 文本"""
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])


def extract_text_pdf(path: str) -> str:
    """解析 PDF 文本"""
    doc = fitz.open(path)
    texts = [page.get_text() for page in doc]
    return "\n".join(texts)


def extract_text_pptx(path: str) -> str:
    """解析 PPT 文本"""
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)


def extract_text_image(path: str) -> str:
    """解析图片信息"""
    return understand_image(path, "请总结图片内容")


def extract_text_audio_vedio(path: str, model_size="small") -> str:
    """
    使用 Whisper 转录音频为文本
    model_size 可选：tiny / base / small / medium / large
    """
    model = whisper.load_model(model_size)
    result = model.transcribe(path)
    return result.get("text", "").strip()

